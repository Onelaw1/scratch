import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation(md_file_path, output_path):
    prs = Presentation()
    
    # Define colors
    PRIMARY_COLOR = RGBColor(0, 51, 102)  # Dark Blue
    SECONDARY_COLOR = RGBColor(255, 255, 255) # White
    ACCENT_COLOR = RGBColor(0, 102, 204) # Lighter Blue

    def add_title_slide(title, subtitle):
        slide_layout = prs.slide_layouts[0] # Title Slide
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        
        title_shape.text = title
        subtitle_shape.text = subtitle
        
        # Style
        title_shape.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
        title_shape.text_frame.paragraphs[0].font.bold = True

    def add_section_slide(title):
        slide_layout = prs.slide_layouts[2] # Section Header (usually) or Title Only
        # Let's use Blank and build it for better control or Title Only
        slide_layout = prs.slide_layouts[6] # Blank
        slide = prs.slides.add_slide(slide_layout)
        
        # Background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = PRIMARY_COLOR
        
        # Text
        left = Inches(1)
        top = Inches(3)
        width = Inches(8)
        height = Inches(1.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = title
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].font.color.rgb = SECONDARY_COLOR
        tf.paragraphs[0].font.size = Pt(44)
        tf.paragraphs[0].font.bold = True

    def add_content_slide(title, content_lines):
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
        title_shape.text_frame.paragraphs[0].font.size = Pt(32)
        title_shape.text_frame.paragraphs[0].font.bold = True
        
        # Content
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.word_wrap = True
        
        for line in content_lines:
            p = tf.add_paragraph()
            clean_line = line.strip().lstrip('- ').lstrip('* ')
            p.text = clean_line
            p.font.size = Pt(18)
            if line.startswith('    '):
                p.level = 1
            else:
                p.level = 0

    # Read Markdown
    with open(md_file_path, 'r', encoding='utf-8') as f:
        lines = f.readlines()

    current_section = ""
    current_slide_title = ""
    current_content = []
    
    # Initial Title Slide
    add_title_slide("국가문화유산진흥원 인사평가 제도 개선 및 운영", "제안서")

    for line in lines:
        line = line.strip()
        if not line:
            continue
            
        if line.startswith('## '): # Section Header
            # Flush previous slide if exists
            if current_slide_title:
                add_content_slide(current_slide_title, current_content)
                current_slide_title = ""
                current_content = []
            
            section_title = line.replace('## ', '')
            add_section_slide(section_title)
            
        elif line.startswith('### '): # Slide Title
            # Flush previous slide
            if current_slide_title:
                add_content_slide(current_slide_title, current_content)
                current_content = []
            
            current_slide_title = line.replace('### ', '')
            
        elif line.startswith('#### '): # Sub-slide or bold text
             # Treat as bold bullet for now or sub-header
             current_content.append(line.replace('#### ', ''))

        elif line.startswith('*') or line.startswith('-'): # Bullet
            current_content.append(line)
        
        elif line.startswith('>'): # Blockquote (Wireframe desc)
            current_content.append(line)
            
        else: # Normal text
            if current_slide_title: # Only add if we are in a slide
                current_content.append(line)

    # Flush last slide
    if current_slide_title:
        add_content_slide(current_slide_title, current_content)

    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

if __name__ == "__main__":
    md_path = r"c:\Users\Administrator\Downloads\Root\K_Heritage_HR_Project\proposal_master.md"
    pptx_path = r"c:\Users\Administrator\Downloads\Root\K_Heritage_HR_Project\proposal_presentation.pptx"
    create_presentation(md_path, pptx_path)
