from pptx import Presentation
from src.services.diagram_service import DiagramService
import os

def test_frameworks():
    os.makedirs("output", exist_ok=True)
    prs = Presentation()
    service = DiagramService()

    # Test 9-Box Grid
    print("Testing 9-Box Grid...")
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank
    data_9box = {
        "Stars": ["Alice", "Bob"],
        "High Potential": ["Charlie"],
        "Enigma": ["David"],
        "Core Players": ["Eve"],
        "Solid Performers": ["Frank"],
        "Inconsistent": ["Grace"],
        "Trusted Professionals": ["Hank"],
        "Effective": ["Ivy"],
        "Low Performers": ["Jack"]
    }
    service.generate_nine_box(prs, slide, data_9box)

    # Test FTE Model
    print("Testing FTE Model...")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    data_fte = {
        "Engineering": {"current": 10.5, "required": 15.0},
        "Sales": {"current": 8.0, "required": 8.0},
        "HR": {"current": 5.0, "required": 4.0},
        "Marketing": {"current": 6.0, "required": 7.5}
    }
    service.generate_fte_model(prs, slide, data_fte)

    # Test AHP
    print("Testing AHP...")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    data_ahp = {
        "goal": "Select Vendor",
        "criteria": {"Cost": 0.4, "Quality": 0.3, "Speed": 0.3},
        "alternatives": {"Vendor A": 0.5, "Vendor B": 0.3, "Vendor C": 0.2}
    }
    service.generate_ahp(prs, slide, data_ahp)

    # Test OCAI
    print("Testing OCAI...")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    data_ocai = {
        "Clan": 0.3,
        "Adhocracy": 0.2,
        "Market": 0.4,
        "Hierarchy": 0.1
    }
    service.generate_ocai(prs, slide, data_ocai)

    prs.save("output/verify_phase9.pptx")
    print("Verification complete. Saved to output/verify_phase9.pptx")

if __name__ == "__main__":
    test_frameworks()
