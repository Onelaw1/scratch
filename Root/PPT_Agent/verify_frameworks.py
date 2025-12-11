from pptx import Presentation
from src.services.diagram_service import DiagramService
import os

def verify_frameworks():
    print("Starting Framework Verification...")
    
    # Initialize
    output_dir = "output/verification"
    os.makedirs(output_dir, exist_ok=True)
    
    service = DiagramService(output_dir)
    prs = Presentation()
    
    # 1. 3C Analysis
    print("Testing 3C Analysis...")
    slide = prs.slides.add_slide(prs.slide_layouts[5]) # Title Only
    slide.shapes.title.text = "3C Analysis"
    data_3c = {
        "Customer": ["High Growth", "Young Demographic", "Mobile First"],
        "Company": ["Strong Brand", "High Quality", "Global Network"],
        "Competitor": ["Low Price", "Fast Delivery", "Local Focus"]
    }
    service.generate_3c_analysis(prs, slide, data_3c)
    
    # 2. Impact/Effort Matrix
    print("Testing Impact/Effort Matrix...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Impact/Effort Matrix"
    data_ie = {
        "Quick Wins": ["UI Update", "Bug Fixes"],
        "Major Projects": ["New Backend", "AI Integration"],
        "Fill-ins": ["Icon Refresh", "Text Update"],
        "Thankless Tasks": ["Legacy Refactor", "Manual Data Entry"]
    }
    service.generate_impact_effort_matrix(prs, slide, data_ie)
    
    # 3. Gantt Chart
    print("Testing Gantt Chart...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Gantt Chart"
    data_gantt = {
        "Phase 1: Planning": {"start": 0.0, "duration": 0.2},
        "Phase 2: Design": {"start": 0.2, "duration": 0.3},
        "Phase 3: Dev": {"start": 0.5, "duration": 0.4},
        "Phase 4: Testing": {"start": 0.9, "duration": 0.1}
    }
    service.generate_gantt_chart(prs, slide, data_gantt)
    
    # 4. Fishbone Diagram
    print("Testing Fishbone Diagram...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Fishbone Diagram"
    data_fishbone = {
        "Problem": "Low Sales",
        "Man": ["Lack of Training", "Low Motivation"],
        "Machine": ["Outdated System", "Slow Network"],
        "Method": ["Inefficient Process", "No Standard"],
        "Material": ["Poor Quality", "Shortage"],
        "Measurement": ["Wrong Metrics", "No Feedback"],
        "Environment": ["High Noise", "Poor Lighting"]
    }
    service.generate_fishbone(prs, slide, data_fishbone)
    
    # 5. Organizational Chart
    print("Testing Organizational Chart...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Organizational Chart"
    data_org = {
        "CEO": "John Doe",
        "C-Level": ["Jane Smith (CFO)", "Bob Wilson (CTO)", "Alice Brown (CMO)"]
    }
    service.generate_org_chart(prs, slide, data_org)
    
    # 6. RACI Matrix
    print("Testing RACI Matrix...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "RACI Matrix"
    data_raci = {
        "Task 1": {"PM": "R", "Dev": "C", "QA": "I", "Ops": "A"},
        "Task 2": {"PM": "A", "Dev": "R", "QA": "C", "Ops": "I"},
        "Task 3": {"PM": "I", "Dev": "A", "QA": "R", "Ops": "C"},
        "Task 4": {"PM": "C", "Dev": "I", "QA": "A", "Ops": "R"}
    }
    service.generate_raci_matrix(prs, slide, data_raci)
    
    # 7. Decision Tree
    print("Testing Decision Tree...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Decision Tree"
    data_tree = {
        "root": "Launch New Product?",
        "options": ["Yes (High Risk)", "No (Low Growth)"]
    }
    service.generate_decision_tree(prs, slide, data_tree)
    
    # 8. Lean Canvas
    print("Testing Lean Canvas...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Lean Canvas"
    data_lean = {
        "Problem": ["Problem A", "Problem B"],
        "Solution": ["Solution A", "Solution B"],
        "Key Metrics": ["Metric A", "Metric B"],
        "Unfair Advantage": ["Advantage A", "Advantage B"],
        "Unique Value Proposition": ["UVP A", "UVP B"],
        "Channels": ["Channel A", "Channel B"],
        "Customer Segments": ["Segment A", "Segment B"],
        "Cost Structure": ["Cost A", "Cost B"],
        "Revenue Streams": ["Revenue A", "Revenue B"]
    }
    service.generate_lean_canvas(prs, slide, data_lean)
    
    # 9. 9-Box Grid
    print("Testing 9-Box Grid...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "9-Box Grid"
    data_9box = {
        "Stars": ["Emp A", "Emp B"],
        "Core Players": ["Emp C", "Emp D"],
        "High Potential": ["Emp E", "Emp F"]
    }
    service.generate_nine_box(prs, slide, data_9box)
    
    # 10. FTE Model
    print("Testing FTE Model...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "FTE Model"
    data_fte = {
        "Dept A": {"current": 5.0, "required": 7.0},
        "Dept B": {"current": 8.0, "required": 6.0},
        "Dept C": {"current": 4.0, "required": 4.0}
    }
    service.generate_fte_model(prs, slide, data_fte)
    
    # 11. AHP
    print("Testing AHP...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "AHP Hierarchy"
    data_ahp = {
        "goal": "Select Vendor",
        "criteria": {"Cost": 0.4, "Quality": 0.3, "Speed": 0.3},
        "alternatives": {"Vendor A": 0.6, "Vendor B": 0.4}
    }
    service.generate_ahp(prs, slide, data_ahp)
    
    # 12. OCAI
    print("Testing OCAI...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "OCAI Culture"
    data_ocai = {
        "Clan": 0.3,
        "Adhocracy": 0.4,
        "Market": 0.2,
        "Hierarchy": 0.1
    }
    service.generate_ocai(prs, slide, data_ocai)
    
    # 13. Futures Wheel
    print("Testing Futures Wheel...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Futures Wheel"
    data_futures = {
        "event": "AI Adoption",
        "consequences": ["Efficiency Up", "Job Changes", "New Risks", "Data Privacy", "Innovation", "Cost Reduction"]
    }
    service.generate_futures_wheel(prs, slide, data_futures)
    
    # 14. Kano Model
    print("Testing Kano Model...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Kano Model"
    data_kano = {
        "Basic": ["Login", "Save"],
        "Performance": ["Speed", "Accuracy"],
        "Delighters": ["AI Insights", "Voice Control"]
    }
    service.generate_kano_model(prs, slide, data_kano)
    
    # Save
    output_path = os.path.join(output_dir, "framework_verification.pptx")
    prs.save(output_path)
    print(f"Verification complete. Saved to {output_path}")

if __name__ == "__main__":
    verify_frameworks()
