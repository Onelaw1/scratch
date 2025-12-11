import os
import sys
from pptx import Presentation

def verify_pptx_content(file_path):
    if not os.path.exists(file_path):
        print(f"Error: File not found at {file_path}")
        return

    prs = Presentation(file_path)
    
    print(f"# PPTX Verification: {os.path.basename(file_path)}\n")
    print(f"**Total Slides:** {len(prs.slides)}\n")
    
    for i, slide in enumerate(prs.slides):
        print(f"## Slide {i+1}")
        
        # Title
        if slide.shapes.title:
            print(f"**Title:** {slide.shapes.title.text}")
        else:
            print("**Title:** (No Title)")
            
        # Text Content
        print("**Content:**")
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if shape == slide.shapes.title:
                continue
                
            for paragraph in shape.text_frame.paragraphs:
                if paragraph.text.strip():
                    prefix = "- " if paragraph.level > 0 else ""
                    print(f"{prefix}{paragraph.text}")
        
        # Shapes/Diagrams (Basic detection)
        print("\n**Visual Elements:**")
        shape_types = []
        for shape in slide.shapes:
            shape_types.append(f"{shape.shape_type} ({shape.name})")
        
        if shape_types:
            for st in shape_types:
                print(f"- {st}")
        else:
            print("(None)")
            
        print("\n---\n")

if __name__ == "__main__":
    if len(sys.argv) > 1:
        file_path = sys.argv[1]
    else:
        # Default to the file generated in the previous step
        file_path = os.path.join(os.path.dirname(__file__), "..", "output", "tests", "real_data_test.pptx")
    
    verify_pptx_content(os.path.abspath(file_path))
