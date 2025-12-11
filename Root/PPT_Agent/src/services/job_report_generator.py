"""
Job Report Generator Service
Generates PowerPoint presentations from Job Management System data.
"""
from pptx import Presentation
from typing import Dict, List, Optional
from .job_data_extractor import JobDataExtractor
from .diagram_service import DiagramService
from .pptx_service import PPTXService
from ..models.schema import PresentationSchema, SlideContent


class JobReportGenerator:
    """Generate job analysis reports as PowerPoint presentations"""
    
    def __init__(self, db_path: str = None, output_dir: str = "output/job_reports"):
        """
        Initialize report generator
        
        Args:
            db_path: Path to Job Management System database
            output_dir: Directory for output files
        """
        self.extractor = JobDataExtractor(db_path)
        self.diagram_service = DiagramService()
        self.pptx_service = PPTXService(output_dir)
        self.output_dir = output_dir
    
    def generate_org_chart(self, prs: Presentation, slide) -> None:
        """
        Generate organizational chart slide
        
        Args:
            prs: Presentation object
            slide: Slide to add chart to
        """
        org_data = self.extractor.get_org_hierarchy()
        self.diagram_service.generate_org_chart(prs, slide, org_data)
    
    def generate_fte_analysis(self, prs: Presentation, slide) -> None:
        """
        Generate FTE analysis slide
        
        Args:
            prs: Presentation object
            slide: Slide to add analysis to
        """
        fte_data = self.extractor.get_fte_by_department()
        self.diagram_service.generate_fte_model(prs, slide, fte_data)
    
    def generate_raci_matrix(self, prs: Presentation, slide) -> None:
        """
        Generate RACI matrix slide
        
        Args:
            prs: Presentation object
            slide: Slide to add matrix to
        """
        raci_data = self.extractor.get_job_tasks_matrix()
        self.diagram_service.generate_raci_matrix(prs, slide, raci_data)
    
    def generate_workload_dashboard(self, prs: Presentation, slide) -> None:
        """
        Generate workload analysis dashboard
        
        Args:
            prs: Presentation object
            slide: Slide to add dashboard to
        """
        workload_stats = self.extractor.get_workload_stats()
        
        # Add title
        title_shape = slide.shapes.title
        title_shape.text = "Workload Analysis Dashboard"
        
        # Add statistics as text boxes
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        
        # Total positions
        textbox = slide.shapes.add_textbox(
            Inches(1), Inches(2), Inches(8), Inches(1)
        )
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        p.text = f"Total Positions: {workload_stats.get('total_positions', 0)}"
        p.font.size = Pt(24)
        p.font.bold = True
        
        # Grade distribution
        y_pos = 3.5
        for grade, count in workload_stats.get('by_grade', {}).items():
            textbox = slide.shapes.add_textbox(
                Inches(1), Inches(y_pos), Inches(8), Inches(0.5)
            )
            tf = textbox.text_frame
            p = tf.paragraphs[0]
            p.text = f"{grade}: {count} positions"
            p.font.size = Pt(16)
            y_pos += 0.6
    
    def generate_grade_analysis(self, prs: Presentation, slide) -> None:
        """
        Generate grade distribution analysis
        
        Args:
            prs: Presentation object
            slide: Slide to add analysis to
        """
        grade_dist = self.extractor.get_grade_distribution()
        
        # Use 9-box grid to visualize grade distribution
        # Transform data for 9-box format
        nine_box_data = {}
        for grade, count in grade_dist.items():
            category = f"Grade {grade}"
            nine_box_data[category] = [f"{count} positions"]
        
        self.diagram_service.generate_nine_box(prs, slide, nine_box_data)
    
    def generate_full_report(
        self, 
        filename: str = "job_analysis_report.pptx",
        language: str = "ko"
    ) -> str:
        """
        Generate complete job analysis report
        
        Args:
            filename: Output filename
            language: Report language ('ko' or 'en')
        
        Returns:
            Path to generated file
        """
        # Create presentation
        prs = Presentation()
        
        # Slide 1: Title
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = title_slide.shapes.title
        subtitle = title_slide.placeholders[1]
        
        if language == "ko":
            title.text = "직무 분석 보고서"
            subtitle.text = "Job Analysis Report"
        else:
            title.text = "Job Analysis Report"
            subtitle.text = "Comprehensive Analysis"
        
        # Slide 2: Organizational Chart
        org_slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title only layout
        title = org_slide.shapes.title
        title.text = "조직도" if language == "ko" else "Organizational Chart"
        self.generate_org_chart(prs, org_slide)
        
        # Slide 3: FTE Analysis
        fte_slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = fte_slide.shapes.title
        title.text = "FTE 분석" if language == "ko" else "FTE Analysis"
        self.generate_fte_analysis(prs, fte_slide)
        
        # Slide 4: RACI Matrix
        raci_slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = raci_slide.shapes.title
        title.text = "RACI 매트릭스" if language == "ko" else "RACI Matrix"
        self.generate_raci_matrix(prs, raci_slide)
        
        # Slide 5: Workload Dashboard
        workload_slide = prs.slides.add_slide(prs.slide_layouts[5])
        self.generate_workload_dashboard(prs, workload_slide)
        
        # Slide 6: Grade Analysis
        grade_slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = grade_slide.shapes.title
        title.text = "등급 분석" if language == "ko" else "Grade Analysis"
        self.generate_grade_analysis(prs, grade_slide)
        
        # Save presentation
        import os
        os.makedirs(self.output_dir, exist_ok=True)
        output_path = os.path.join(self.output_dir, filename)
        prs.save(output_path)
        
        return output_path
    
    def close(self):
        """Clean up resources"""
        self.extractor.close()
