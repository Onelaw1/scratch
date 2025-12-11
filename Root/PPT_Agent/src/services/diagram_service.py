from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from typing import Dict, List
import os

class DiagramService:
    def __init__(self, output_dir: str = "output/diagrams"):
        self.output_dir = output_dir
        os.makedirs(self.output_dir, exist_ok=True)

    def generate_swot(self, prs: Presentation, slide, data: Dict[str, List[str]]) -> None:
        """Generate SWOT Analysis diagram directly on slide."""
        print("DiagramService: Generating SWOT Analysis...")
        
        # Define quadrant positions
        left_col = Inches(0.5)
        right_col = Inches(5.5)
        top_row = Inches(1.5)
        bottom_row = Inches(4.2)
        width = Inches(4.5)
        height = Inches(2.5)
        
        # Colors
        colors = {
            "Strengths": RGBColor(46, 134, 171),    # Blue
            "Weaknesses": RGBColor(162, 59, 114),   # Purple
            "Opportunities": RGBColor(241, 143, 1), # Orange
            "Threats": RGBColor(199, 62, 29)        # Red
        }
        
        quadrants = [
            ("Strengths", left_col, top_row),
            ("Weaknesses", right_col, top_row),
            ("Opportunities", left_col, bottom_row),
            ("Threats", right_col, bottom_row)
        ]
        
        for title, left, top in quadrants:
            # Box
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left, top, width, height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = colors[title]
            shape.line.color.rgb = RGBColor(255, 255, 255)
            shape.line.width = Pt(2)
            
            # Text
            text_frame = shape.text_frame
            text_frame.margin_top = Inches(0.1)
            text_frame.margin_left = Inches(0.2)
            
            # Title
            p = text_frame.paragraphs[0]
            p.text = title
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            
            # Bullet points
            for item in data.get(title, [])[:3]:  # Max 3 items
                p = text_frame.add_paragraph()
                p.text = f"• {item}"
                p.font.size = Pt(12)
                p.font.color.rgb = RGBColor(255, 255, 255)
                p.level = 0

    def generate_workflow(self, prs: Presentation, slide, steps: List[str]) -> None:
        """Generate horizontal workflow diagram."""
        print("DiagramService: Generating Workflow...")
        
        num_steps = min(len(steps), 5)  # Max 5 steps
        box_width = Inches(1.8)
        box_height = Inches(1.2)
        spacing = Inches(0.3)
        start_left = Inches(0.5)
        top = Inches(2.5)
        
        for i, step in enumerate(steps[:num_steps]):
            left = start_left + i * (box_width + spacing)
            
            # Box
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left, top, box_width, box_height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(46, 134, 171)
            shape.line.color.rgb = RGBColor(30, 90, 120)
            shape.line.width = Pt(2)
            
            # Text
            text_frame = shape.text_frame
            text_frame.word_wrap = True
            text_frame.vertical_anchor = 1  # Middle
            
            p = text_frame.paragraphs[0]
            p.text = f"{i+1}. {step}"
            p.font.size = Pt(11)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER
            
            # Arrow (if not last step)
            if i < num_steps - 1:
                arrow_left = left + box_width
                arrow = slide.shapes.add_shape(
                    MSO_SHAPE.RIGHT_ARROW,
                    arrow_left, top + Inches(0.4),
                    spacing, Inches(0.4)
                )
                arrow.fill.solid()
                arrow.fill.fore_color.rgb = RGBColor(100, 100, 100)
                arrow.line.fill.background()

    def generate_four_p(self, prs: Presentation, slide, data: Dict[str, List[str]]) -> None:
        """Generate Marketing 4P Framework."""
        print("DiagramService: Generating 4P Framework...")
        
        # Center circle
        center_x = Inches(5)
        center_y = Inches(3.5)
        circle_size = Inches(1.5)
        
        center_circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            center_x - circle_size/2, center_y - circle_size/2,
            circle_size, circle_size
        )
        center_circle.fill.solid()
        center_circle.fill.fore_color.rgb = RGBColor(46, 134, 171)
        
        tf = center_circle.text_frame
        p = tf.paragraphs[0]
        p.text = "4P\nMarketing"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        # Four quadrants
        ps = ["Product", "Price", "Place", "Promotion"]
        positions = [
            (Inches(1), Inches(1.5)),    # Top-left
            (Inches(7.5), Inches(1.5)),  # Top-right
            (Inches(1), Inches(5)),      # Bottom-left
            (Inches(7.5), Inches(5))     # Bottom-right
        ]
        colors = [
            RGBColor(241, 143, 1),
            RGBColor(162, 59, 114),
            RGBColor(106, 153, 78),
            RGBColor(199, 62, 29)
        ]
        
        for i, (p_name, (left, top)) in enumerate(zip(ps, positions)):
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left, top, Inches(2), Inches(1.5)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = colors[i]
            
            tf = shape.text_frame
            tf.margin_top = Inches(0.1)
            tf.margin_left = Inches(0.1)
            
            p = tf.paragraphs[0]
            p.text = p_name
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            
            for item in data.get(p_name, [])[:2]:
                p = tf.add_paragraph()
                p.text = f"• {item}"
                p.font.size = Pt(10)
                p.font.color.rgb = RGBColor(255, 255, 255)

    def generate_bcg_matrix(self, prs: Presentation, slide, data: Dict[str, List[str]]) -> None:
        """Generate BCG Matrix (2x2)."""
        print("DiagramService: Generating BCG Matrix...")
        
        # Define quadrant positions
        left_col = Inches(1)
        right_col = Inches(5.5)
        top_row = Inches(2)
        bottom_row = Inches(4.5)
        width = Inches(4)
        height = Inches(2.2)
        
        # Colors
        colors = {
            "Stars": RGBColor(241, 143, 1),        # Orange
            "Cash Cows": RGBColor(106, 153, 78),   # Green
            "Question Marks": RGBColor(162, 59, 114), # Purple
            "Dogs": RGBColor(150, 150, 150)        # Grey
        }
        
        quadrants = [
            ("Stars", right_col, top_row, "High Growth\nHigh Share"),
            ("Question Marks", left_col, top_row, "High Growth\nLow Share"),
            ("Cash Cows", right_col, bottom_row, "Low Growth\nHigh Share"),
            ("Dogs", left_col, bottom_row, "Low Growth\nLow Share")
        ]
        
        for title, left, top, subtitle in quadrants:
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left, top, width, height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = colors[title]
            shape.line.color.rgb = RGBColor(255, 255, 255)
            shape.line.width = Pt(2)
            
            tf = shape.text_frame
            tf.margin_top = Inches(0.1)
            tf.margin_left = Inches(0.15)
            
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            
            p = tf.add_paragraph()
            p.text = subtitle
            p.font.size = Pt(10)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.italic = True
            
            for item in data.get(title, [])[:2]:
                p = tf.add_paragraph()
                p.text = f"• {item}"
                p.font.size = Pt(11)
                p.font.color.rgb = RGBColor(255, 255, 255)

    def generate_eisenhower_matrix(self, prs: Presentation, slide, data: Dict[str, List[str]]) -> None:
        """Generate Eisenhower Matrix (Urgent/Important)."""
        print("DiagramService: Generating Eisenhower Matrix...")
        
        left_col = Inches(1)
        right_col = Inches(5.5)
        top_row = Inches(2)
        bottom_row = Inches(4.5)
        width = Inches(4)
        height = Inches(2.2)
        
        colors = {
            "Do First": RGBColor(199, 62, 29),     # Red
            "Schedule": RGBColor(241, 143, 1),     # Orange
            "Delegate": RGBColor(46, 134, 171),    # Blue
            "Eliminate": RGBColor(150, 150, 150)   # Grey
        }
        
        quadrants = [
            ("Do First", right_col, top_row, "Urgent & Important"),
            ("Schedule", left_col, top_row, "Not Urgent\nBut Important"),
            ("Delegate", right_col, bottom_row, "Urgent\nNot Important"),
            ("Eliminate", left_col, bottom_row, "Not Urgent\nNot Important")
        ]
        
        for title, left, top, subtitle in quadrants:
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left, top, width, height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = colors[title]
            shape.line.color.rgb = RGBColor(255, 255, 255)
            shape.line.width = Pt(2)
            
            tf = shape.text_frame
            tf.margin_top = Inches(0.1)
            tf.margin_left = Inches(0.15)
            
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            
            p = tf.add_paragraph()
            p.text = subtitle
            p.font.size = Pt(10)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.italic = True
            
            for item in data.get(title, [])[:2]:
                p = tf.add_paragraph()
                p.text = f"• {item}"
                p.font.size = Pt(11)
                p.font.color.rgb = RGBColor(255, 255, 255)

    def generate_porters_five_forces(self, prs: Presentation, slide, data: Dict[str, str]) -> None:
        """Generate Porter's 5 Forces diagram."""
        print("DiagramService: Generating Porter's 5 Forces...")
        
        # Center box
        center_x = Inches(3.5)
        center_y = Inches(3.5)
        center_w = Inches(3)
        center_h = Inches(1.5)
        
        center = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            center_x, center_y, center_w, center_h
        )
        center.fill.solid()
        center.fill.fore_color.rgb = RGBColor(46, 134, 171)
        center.line.width = Pt(3)
        
        tf = center.text_frame
        p = tf.paragraphs[0]
        p.text = "Industry\nRivalry"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        # Five forces positions
        forces = [
            ("New Entrants", Inches(3.5), Inches(0.5), "top"),
            ("Suppliers", Inches(0.3), Inches(3.5), "left"),
            ("Buyers", Inches(6.7), Inches(3.5), "right"),
            ("Substitutes", Inches(3.5), Inches(6.2), "bottom")
        ]
        
        force_w = Inches(2.5)
        force_h = Inches(1.2)
        
        for name, left, top, position in forces:
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left, top, force_w, force_h
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(241, 143, 1)
            shape.line.width = Pt(2)
            
            tf = shape.text_frame
            p = tf.paragraphs[0]
            p.text = name
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER
            
            # Add description if provided
            if name in data:
                p = tf.add_paragraph()
                p.text = data[name]
                p.font.size = Pt(9)
                p.font.color.rgb = RGBColor(255, 255, 255)
                p.alignment = PP_ALIGN.CENTER

    def generate_ge_matrix(self, prs: Presentation, slide, data: Dict[str, List[str]]) -> None:
        """Generate GE/McKinsey Matrix (3x3)."""
        print("DiagramService: Generating GE Matrix...")
        
        # 3x3 grid positions
        cell_width = Inches(3)
        cell_height = Inches(1.8)
        start_left = Inches(1)
        start_top = Inches(1.5)
        
        # Define 9 cells with colors
        cells = [
            # Row 1 (High Industry Attractiveness)
            ("High/High", 0, 0, RGBColor(106, 153, 78), "Invest/Grow"),
            ("Med/High", 1, 0, RGBColor(106, 153, 78), "Invest/Grow"),
            ("Low/High", 2, 0, RGBColor(241, 143, 1), "Selective"),
            # Row 2 (Medium Industry Attractiveness)
            ("High/Med", 0, 1, RGBColor(106, 153, 78), "Invest/Grow"),
            ("Med/Med", 1, 1, RGBColor(241, 143, 1), "Selective"),
            ("Low/Med", 2, 1, RGBColor(199, 62, 29), "Harvest"),
            # Row 3 (Low Industry Attractiveness)
            ("High/Low", 0, 2, RGBColor(241, 143, 1), "Selective"),
            ("Med/Low", 1, 2, RGBColor(199, 62, 29), "Harvest"),
            ("Low/Low", 2, 2, RGBColor(199, 62, 29), "Divest")
        ]
        
        for cell_name, col, row, color, strategy in cells:
            left = start_left + col * cell_width
            top = start_top + row * cell_height
            
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                left, top, cell_width, cell_height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = color
            shape.line.color.rgb = RGBColor(255, 255, 255)
            shape.line.width = Pt(2)
            
            tf = shape.text_frame
            tf.margin_top = Inches(0.05)
            tf.margin_left = Inches(0.1)
            
            p = tf.paragraphs[0]
            p.text = strategy
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            
            # Add products in this cell
            if cell_name in data:
                for item in data.get(cell_name, [])[:2]:
                    p = tf.add_paragraph()
                    p.text = f"• {item}"
                    p.font.size = Pt(9)
                    p.font.color.rgb = RGBColor(255, 255, 255)

    def generate_scenario_planning(self, prs: Presentation, slide, data: Dict[str, Dict]) -> None:
        """Generate Scenario Planning 2x2 matrix."""
        print("DiagramService: Generating Scenario Planning...")
        
        left_col = Inches(1)
        right_col = Inches(5.5)
        top_row = Inches(2)
        bottom_row = Inches(4.5)
        width = Inches(4)
        height = Inches(2.2)
        
        # Define scenarios
        scenarios = [
            ("Scenario 1", right_col, top_row, RGBColor(106, 153, 78)),
            ("Scenario 2", left_col, top_row, RGBColor(241, 143, 1)),
            ("Scenario 3", right_col, bottom_row, RGBColor(162, 59, 114)),
            ("Scenario 4", left_col, bottom_row, RGBColor(199, 62, 29))
        ]
        
        for scenario_name, left, top, color in scenarios:
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left, top, width, height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = color
            shape.line.color.rgb = RGBColor(255, 255, 255)
            shape.line.width = Pt(2)
            
            tf = shape.text_frame
            tf.margin_top = Inches(0.1)
            tf.margin_left = Inches(0.15)
            
            scenario_data = data.get(scenario_name, {})
            
            p = tf.paragraphs[0]
            p.text = scenario_data.get("name", scenario_name)
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            
            # Description
            if "description" in scenario_data:
                p = tf.add_paragraph()
                p.text = scenario_data["description"]
                p.font.size = Pt(10)
                p.font.color.rgb = RGBColor(255, 255, 255)
            
            # Key factors
            for factor in scenario_data.get("factors", [])[:2]:
                p = tf.add_paragraph()
                p.text = f"• {factor}"
                p.font.size = Pt(9)
                p.font.color.rgb = RGBColor(255, 255, 255)

    def generate_value_chain(self, prs: Presentation, slide, data: Dict[str, List[str]]) -> None:
        """Generate Porter's Value Chain."""
        print("DiagramService: Generating Value Chain...")
        
        # Primary activities (bottom row)
        primary_width = Inches(1.8)
        primary_height = Inches(1.5)
        primary_top = Inches(4.5)
        primary_start = Inches(0.5)
        
        primary_activities = ["Inbound", "Operations", "Outbound", "Marketing", "Service"]
        primary_color = RGBColor(46, 134, 171)
        
        for i, activity in enumerate(primary_activities):
            left = primary_start + i * primary_width
            
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left, primary_top, primary_width, primary_height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = primary_color
            shape.line.color.rgb = RGBColor(255, 255, 255)
            shape.line.width = Pt(2)
            
            tf = shape.text_frame
            tf.margin_top = Inches(0.05)
            tf.margin_left = Inches(0.1)
            
            p = tf.paragraphs[0]
            p.text = activity
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER
            
            # Add items
            for item in data.get(activity, [])[:1]:
                p = tf.add_paragraph()
                p.text = f"• {item}"
                p.font.size = Pt(8)
                p.font.color.rgb = RGBColor(255, 255, 255)
        
        # Support activities (top row)
        support_width = Inches(9)
        support_height = Inches(0.7)
        support_start = Inches(0.5)
        support_activities = ["Infrastructure", "HR Management", "Technology", "Procurement"]
        support_color = RGBColor(241, 143, 1)
        
        for i, activity in enumerate(support_activities):
            top = Inches(1.5) + i * support_height
            
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                support_start, top, support_width, support_height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = support_color
            shape.line.color.rgb = RGBColor(255, 255, 255)
            shape.line.width = Pt(2)
            
            tf = shape.text_frame
            tf.margin_top = Inches(0.05)
            tf.margin_left = Inches(0.1)
            
            p = tf.paragraphs[0]
            p.text = activity
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)

    def generate_business_model_canvas(self, prs: Presentation, slide, data: Dict[str, List[str]]) -> None:
        """Generate Business Model Canvas (9 blocks)."""
        print("DiagramService: Generating Business Model Canvas...")
        
        # Define 9 blocks with positions
        blocks = [
            # Left column
            ("Key Partners", Inches(0.5), Inches(1.5), Inches(2.2), Inches(4.5), RGBColor(241, 143, 1)),
            ("Key Activities", Inches(2.8), Inches(1.5), Inches(2.2), Inches(2.2), RGBColor(46, 134, 171)),
            ("Key Resources", Inches(2.8), Inches(3.8), Inches(2.2), Inches(2.2), RGBColor(46, 134, 171)),
            # Center column
            ("Value Propositions", Inches(5.1), Inches(1.5), Inches(2.2), Inches(4.5), RGBColor(199, 62, 29)),
            # Right column
            ("Customer Relationships", Inches(7.4), Inches(1.5), Inches(2.2), Inches(2.2), RGBColor(106, 153, 78)),
            ("Channels", Inches(7.4), Inches(3.8), Inches(2.2), Inches(2.2), RGBColor(106, 153, 78)),
            ("Customer Segments", Inches(9.7), Inches(1.5), Inches(2.2), Inches(4.5), RGBColor(162, 59, 114)),
            # Bottom row
            ("Cost Structure", Inches(0.5), Inches(6.1), Inches(4.6), Inches(1.4), RGBColor(150, 150, 150)),
            ("Revenue Streams", Inches(5.1), Inches(6.1), Inches(6.8), Inches(1.4), RGBColor(106, 153, 78))
        ]
        
        for block_name, left, top, width, height, color in blocks:
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left, top, width, height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = color
            shape.line.color.rgb = RGBColor(255, 255, 255)
            shape.line.width = Pt(2)
            
            tf = shape.text_frame
            tf.margin_top = Inches(0.05)
            tf.margin_left = Inches(0.1)
            tf.word_wrap = True
            
            p = tf.paragraphs[0]
            p.text = block_name
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            
            # Add items
            for item in data.get(block_name, [])[:2]:
                p = tf.add_paragraph()
                p.text = f"• {item}"
                p.font.size = Pt(8)
                p.font.color.rgb = RGBColor(255, 255, 255)

    def generate_mckinsey_7s(self, prs: Presentation, slide, data: Dict[str, str]) -> None:
        """Generate McKinsey 7S Framework."""
        print("DiagramService: Generating McKinsey 7S...")
        
        # Center circle
        center_x = Inches(5)
        center_y = Inches(3.75)
        center_size = Inches(1.5)
        
        center = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            center_x - center_size/2, center_y - center_size/2,
            center_size, center_size
        )
        center.fill.solid()
        center.fill.fore_color.rgb = RGBColor(199, 62, 29)
        center.line.width = Pt(3)
        
        tf = center.text_frame
        p = tf.paragraphs[0]
        p.text = "Shared\nValues"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        # 6 surrounding circles
        import math
        radius = Inches(2.5)
        circle_size = Inches(1.3)
        elements = ["Strategy", "Structure", "Systems", "Style", "Staff", "Skills"]
        colors = [
            RGBColor(46, 134, 171),
            RGBColor(241, 143, 1),
            RGBColor(106, 153, 78),
            RGBColor(162, 59, 114),
            RGBColor(46, 134, 171),
            RGBColor(241, 143, 1)
        ]
        
        for i, (element, color) in enumerate(zip(elements, colors)):
            angle = (i * 60 - 90) * math.pi / 180  # Start from top
            x = center_x + radius * math.cos(angle) - circle_size/2
            y = center_y + radius * math.sin(angle) - circle_size/2
            
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                x, y, circle_size, circle_size
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = color
            circle.line.width = Pt(2)
            
            tf = circle.text_frame
            p = tf.paragraphs[0]
            p.text = element
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER

    def generate_ansoff_matrix(self, prs: Presentation, slide, data: Dict[str, List[str]]) -> None:
        """Generate Ansoff Matrix (2x2)."""
        print("DiagramService: Generating Ansoff Matrix...")
        
        left_col = Inches(1)
        right_col = Inches(5.5)
        top_row = Inches(2)
        bottom_row = Inches(4.5)
        width = Inches(4)
        height = Inches(2.2)
        
        colors = {
            "Market Penetration": RGBColor(106, 153, 78),  # Green - Low risk
            "Product Development": RGBColor(241, 143, 1),  # Orange - Medium risk
            "Market Development": RGBColor(241, 143, 1),   # Orange - Medium risk
            "Diversification": RGBColor(199, 62, 29)       # Red - High risk
        }
        
        quadrants = [
            ("Market Penetration", left_col, top_row, "Existing Market\nExisting Product"),
            ("Product Development", right_col, top_row, "Existing Market\nNew Product"),
            ("Market Development", left_col, bottom_row, "New Market\nExisting Product"),
            ("Diversification", right_col, bottom_row, "New Market\nNew Product")
        ]
        
        for title, left, top, subtitle in quadrants:
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left, top, width, height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = colors[title]
            shape.line.color.rgb = RGBColor(255, 255, 255)
            shape.line.width = Pt(2)
            
            tf = shape.text_frame
            tf.margin_top = Inches(0.1)
            tf.margin_left = Inches(0.15)
            
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            
            p = tf.add_paragraph()
            p.text = subtitle
            p.font.size = Pt(10)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.italic = True
            
            for item in data.get(title, [])[:2]:
                p = tf.add_paragraph()
                p.text = f"• {item}"
                p.font.size = Pt(10)
                p.font.color.rgb = RGBColor(255, 255, 255)

    def generate_stakeholder_mapping(self, prs: Presentation, slide, data: Dict[str, Dict]) -> None:
        """Generate Stakeholder Mapping (Power/Interest Matrix)."""
        print("DiagramService: Generating Stakeholder Mapping...")
        
        # Draw quadrants
        left_col = Inches(1)
        right_col = Inches(5.5)
        top_row = Inches(2)
        bottom_row = Inches(4.5)
        width = Inches(4)
        height = Inches(2.2)
        
        quadrants = [
            ("Manage Closely", right_col, top_row, RGBColor(199, 62, 29)),
            ("Keep Satisfied", left_col, top_row, RGBColor(241, 143, 1)),
            ("Keep Informed", right_col, bottom_row, RGBColor(106, 153, 78)),
            ("Monitor", left_col, bottom_row, RGBColor(150, 150, 150))
        ]
        
        for title, left, top, color in quadrants:
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                left, top, width, height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = color
            shape.fill.transparency = 0.5
            shape.line.color.rgb = RGBColor(255, 255, 255)
            shape.line.width = Pt(2)
            
            tf = shape.text_frame
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER
        
        # Add stakeholder circles
        for stakeholder_name, stakeholder_data in data.items():
            if isinstance(stakeholder_data, dict):
                power = stakeholder_data.get("power", 0.5)  # 0-1
                interest = stakeholder_data.get("interest", 0.5)  # 0-1
                
                # Map to position
                x = Inches(1) + Inches(8) * interest
                y = Inches(6.7) - Inches(4.4) * power
                
                circle = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    x - Inches(0.3), y - Inches(0.3),
                    Inches(0.6), Inches(0.6)
                )
                circle.fill.solid()
                circle.fill.fore_color.rgb = RGBColor(46, 134, 171)
                circle.line.color.rgb = RGBColor(255, 255, 255)
                circle.line.width = Pt(2)
                
                # Add label
                tf = circle.text_frame
                p = tf.paragraphs[0]
                p.text = stakeholder_name[:3]  # Abbreviation
                p.font.size = Pt(8)
                p.font.bold = True
                p.font.color.rgb = RGBColor(255, 255, 255)
                p.alignment = PP_ALIGN.CENTER

    def generate_3c_analysis(self, prs: Presentation, slide, data: Dict[str, List[str]]) -> None:
        """Generate 3C Analysis (Customer, Company, Competitor)."""
        print("DiagramService: Generating 3C Analysis...")
        
        # Triangle arrangement of 3 circles
        import math
        center_x = Inches(5)
        center_y = Inches(4)
        radius = Inches(2)
        circle_size = Inches(2.5)
        
        elements = ["Customer", "Company", "Competitor"]
        colors = [
            RGBColor(46, 134, 171),   # Blue
            RGBColor(241, 143, 1),    # Orange
            RGBColor(106, 153, 78)    # Green
        ]
        
        # Triangle positions (top, bottom-left, bottom-right)
        positions = [
            (center_x - circle_size/2, center_y - radius - circle_size/2),  # Top
            (center_x - radius * 0.866 - circle_size/2, center_y + radius/2 - circle_size/2),  # Bottom-left
            (center_x + radius * 0.866 - circle_size/2, center_y + radius/2 - circle_size/2)   # Bottom-right
        ]
        
        for i, (element, color) in enumerate(zip(elements, colors)):
            x, y = positions[i]
            
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                x, y, circle_size, circle_size
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = color
            circle.line.color.rgb = RGBColor(255, 255, 255)
            circle.line.width = Pt(3)
            
            tf = circle.text_frame
            tf.margin_top = Inches(0.1)
            tf.margin_left = Inches(0.1)
            
            p = tf.paragraphs[0]
            p.text = element
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER
            
            # Add items
            for item in data.get(element, [])[:2]:
                p = tf.add_paragraph()
                p.text = f"• {item}"
                p.font.size = Pt(10)
                p.font.color.rgb = RGBColor(255, 255, 255)
                p.alignment = PP_ALIGN.CENTER
        
        # Add "Strategic Fit" in center
        center_text = slide.shapes.add_textbox(
            center_x - Inches(1), center_y - Inches(0.3),
            Inches(2), Inches(0.6)
        )
        tf = center_text.text_frame
        p = tf.paragraphs[0]
        p.text = "Strategic\nFit"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.alignment = PP_ALIGN.CENTER

    def generate_impact_effort_matrix(self, prs: Presentation, slide, data: Dict[str, List[str]]) -> None:
        """Generate Impact/Effort Matrix (2x2)."""
        print("DiagramService: Generating Impact/Effort Matrix...")
        
        left_col = Inches(1)
        right_col = Inches(5.5)
        top_row = Inches(2)
        bottom_row = Inches(4.5)
        width = Inches(4)
        height = Inches(2.2)
        
        colors = {
            "Quick Wins": RGBColor(106, 153, 78),      # Green - High impact, low effort
            "Major Projects": RGBColor(241, 143, 1),   # Orange - High impact, high effort
            "Fill-ins": RGBColor(150, 150, 150),       # Grey - Low impact, low effort
            "Thankless Tasks": RGBColor(199, 62, 29)   # Red - Low impact, high effort
        }
        
        quadrants = [
            ("Quick Wins", left_col, top_row, "High Impact\nLow Effort"),
            ("Major Projects", right_col, top_row, "High Impact\nHigh Effort"),
            ("Fill-ins", left_col, bottom_row, "Low Impact\nLow Effort"),
            ("Thankless Tasks", right_col, bottom_row, "Low Impact\nHigh Effort")
        ]
        
        for title, left, top, subtitle in quadrants:
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left, top, width, height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = colors[title]
            shape.line.color.rgb = RGBColor(255, 255, 255)
            shape.line.width = Pt(2)
            
            tf = shape.text_frame
            tf.margin_top = Inches(0.1)
            tf.margin_left = Inches(0.15)
            
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            
            p = tf.add_paragraph()
            p.text = subtitle
            p.font.size = Pt(10)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.italic = True
            
            for item in data.get(title, [])[:2]:
                p = tf.add_paragraph()
                p.text = f"• {item}"
                p.font.size = Pt(10)
                p.font.color.rgb = RGBColor(255, 255, 255)

    def generate_gantt_chart(self, prs: Presentation, slide, data: Dict[str, Dict]) -> None:
        """Generate Gantt Chart."""
        print("DiagramService: Generating Gantt Chart...")
        
        # Chart area
        start_left = Inches(2)
        start_top = Inches(2)
        row_height = Inches(0.6)
        timeline_width = Inches(7)
        
        # Draw tasks
        for i, (task_name, task_data) in enumerate(list(data.items())[:5]):  # Max 5 tasks
            top = start_top + i * row_height
            
            # Task label
            label_box = slide.shapes.add_textbox(
                Inches(0.5), top, Inches(1.3), row_height
            )
            tf = label_box.text_frame
            p = tf.paragraphs[0]
            p.text = task_name
            p.font.size = Pt(10)
            p.font.bold = True
            p.alignment = PP_ALIGN.RIGHT
            
            # Task bar
            if isinstance(task_data, dict):
                start = task_data.get("start", 0)  # 0-1 (percentage of timeline)
                duration = task_data.get("duration", 0.2)  # 0-1
                
                bar_left = start_left + timeline_width * start
                bar_width = timeline_width * duration
                
                bar = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    bar_left, top + Inches(0.1),
                    bar_width, row_height - Inches(0.2)
                )
                bar.fill.solid()
                bar.fill.fore_color.rgb = RGBColor(46, 134, 171)
                bar.line.color.rgb = RGBColor(255, 255, 255)
                bar.line.width = Pt(1)

    def generate_fishbone(self, prs: Presentation, slide, data: Dict[str, List[str]]) -> None:
        """Generate Fishbone Diagram (Ishikawa)."""
        print("DiagramService: Generating Fishbone Diagram...")
        
        # Central spine
        spine_start = Inches(1)
        spine_end = Inches(9)
        spine_y = Inches(4)
        
        # Draw main spine (arrow)
        spine = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            spine_start, spine_y, spine_end, spine_y
        )
        spine.line.width = Pt(3)
        spine.line.color.rgb = RGBColor(50, 50, 50)
        
        # Arrow head
        head = slide.shapes.add_shape(
            MSO_SHAPE.ISOSCELES_TRIANGLE,
            spine_end - Inches(0.3), spine_y - Inches(0.2),
            Inches(0.4), Inches(0.4)
        )
        head.fill.solid()
        head.fill.fore_color.rgb = RGBColor(50, 50, 50)
        head.line.width = Pt(0)
        head.rotation = 90
        
        # 6M categories
        categories = ["Man", "Machine", "Method", "Material", "Measurement", "Environment"]
        colors = [
            RGBColor(199, 62, 29), RGBColor(241, 143, 1), RGBColor(106, 153, 78),
            RGBColor(46, 134, 171), RGBColor(162, 59, 114), RGBColor(150, 150, 150)
        ]
        
        # Draw bones (3 on top, 3 on bottom)
        for i, (category, color) in enumerate(zip(categories, colors)):
            is_top = i < 3
            x_pos = spine_start + Inches(1.5 + i % 3 * 2.5)
            y_offset = Inches(1.5) if is_top else Inches(-1.5)
            
            # Diagonal bone
            bone = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                x_pos, spine_y,
                x_pos + Inches(1.2), spine_y + y_offset
            )
            bone.line.width = Pt(2)
            bone.line.color.rgb = color
            
            # Category label
            label = slide.shapes.add_textbox(
                x_pos + Inches(0.8), spine_y + y_offset - Inches(0.3),
                Inches(1.5), Inches(0.5)
            )
            tf = label.text_frame
            p = tf.paragraphs[0]
            p.text = category
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = color
            
            # Add causes
            for j, cause in enumerate(data.get(category, [])[:2]):
                cause_y = spine_y + y_offset + (Inches(0.3) if is_top else Inches(-0.6)) + j * Inches(0.3) * (-1 if is_top else 1)
                cause_box = slide.shapes.add_textbox(
                    x_pos + Inches(1.3), cause_y,
                    Inches(1.2), Inches(0.25)
                )
                tf = cause_box.text_frame
                p = tf.paragraphs[0]
                p.text = f"• {cause}"
                p.font.size = Pt(8)
        
        # Problem/Effect label
        problem_box = slide.shapes.add_textbox(
            spine_end - Inches(1.5), spine_y - Inches(0.8),
            Inches(1.5), Inches(0.5)
        )
        tf = problem_box.text_frame
        p = tf.paragraphs[0]
        p.text = data.get("Problem", "Effect")
        p.font.size = Pt(14)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

    def generate_org_chart(self, prs: Presentation, slide, data: Dict) -> None:
        """Generate Organizational Chart."""
        print("DiagramService: Generating Organizational Chart...")
        
        def draw_node(name, level, position, total_at_level):
            box_width = Inches(1.5)
            box_height = Inches(0.6)
            x = Inches(1) + position * Inches(8) / max(total_at_level, 1)
            y = Inches(1.5) + level * Inches(1.2)
            
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x, y, box_width, box_height
            )
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(46, 134, 171) if level == 0 else RGBColor(241, 143, 1) if level == 1 else RGBColor(106, 153, 78)
            box.line.width = Pt(2)
            
            tf = box.text_frame
            p = tf.paragraphs[0]
            p.text = name
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER
            
            return (x + box_width/2, y + box_height)
        
        # Simple 3-level org chart
        ceo = data.get("CEO", "CEO")
        c_level = data.get("C-Level", ["CFO", "CTO", "CMO"])
        
        # Draw CEO
        ceo_pos = draw_node(ceo, 0, 0.5, 1)
        
        # Draw C-Level
        for i, exec in enumerate(c_level[:3]):
            exec_pos = draw_node(exec, 1, i / 3, 3)
            # Connector
            conn = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                ceo_pos[0], ceo_pos[1],
                exec_pos[0], exec_pos[1] - Inches(0.6)
            )
            conn.line.width = Pt(1.5)

    def generate_raci_matrix(self, prs: Presentation, slide, data: Dict[str, Dict[str, str]]) -> None:
        """Generate RACI Matrix."""
        print("DiagramService: Generating RACI Matrix...")
        
        tasks = list(data.keys())[:4]
        roles = ["PM", "Dev", "QA", "Ops"]
        
        cell_width = Inches(1.8)
        cell_height = Inches(0.6)
        start_x = Inches(1.5)
        start_y = Inches(2)
        
        # Header row (roles)
        for i, role in enumerate(roles):
            box = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                start_x + (i+1) * cell_width, start_y,
                cell_width, cell_height
            )
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(46, 134, 171)
            box.line.width = Pt(1)
            
            tf = box.text_frame
            p = tf.paragraphs[0]
            p.text = role
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER
        
        # Task rows
        raci_colors = {
            "R": RGBColor(199, 62, 29),
            "A": RGBColor(46, 134, 171),
            "C": RGBColor(241, 143, 1),
            "I": RGBColor(150, 150, 150)
        }
        
        for row, task in enumerate(tasks):
            # Task label
            label_box = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                start_x, start_y + (row+1) * cell_height,
                cell_width, cell_height
            )
            label_box.fill.solid()
            label_box.fill.fore_color.rgb = RGBColor(200, 200, 200)
            label_box.line.width = Pt(1)
            
            tf = label_box.text_frame
            p = tf.paragraphs[0]
            p.text = task
            p.font.size = Pt(9)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
            
            # RACI cells
            task_data = data.get(task, {})
            for col, role in enumerate(roles):
                raci_value = task_data.get(role, "")
                
                cell = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    start_x + (col+1) * cell_width, start_y + (row+1) * cell_height,
                    cell_width, cell_height
                )
                cell.fill.solid()
                cell.fill.fore_color.rgb = raci_colors.get(raci_value, RGBColor(255, 255, 255))
                cell.line.width = Pt(1)
                
                tf = cell.text_frame
                p = tf.paragraphs[0]
                p.text = raci_value
                p.font.size = Pt(14)
                p.font.bold = True
                p.font.color.rgb = RGBColor(255, 255, 255) if raci_value else RGBColor(0, 0, 0)
                p.alignment = PP_ALIGN.CENTER

    def generate_decision_tree(self, prs: Presentation, slide, data: Dict) -> None:
        """Generate Decision Tree."""
        print("DiagramService: Generating Decision Tree...")
        
        # Simple 2-level decision tree
        root_x = Inches(5)
        root_y = Inches(2)
        
        # Root decision
        root = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            root_x - Inches(0.8), root_y,
            Inches(1.6), Inches(0.6)
        )
        root.fill.solid()
        root.fill.fore_color.rgb = RGBColor(46, 134, 171)
        root.line.width = Pt(2)
        
        tf = root.text_frame
        p = tf.paragraphs[0]
        p.text = data.get("root", "Decision")
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        # Branches
        options = data.get("options", ["Option A", "Option B"])
        for i, option in enumerate(options[:2]):
            branch_x = Inches(2 + i * 6)
            branch_y = Inches(4)
            
            # Connector
            conn = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                root_x, root_y + Inches(0.6),
                branch_x, branch_y
            )
            conn.line.width = Pt(1.5)
            
            # Outcome
            outcome = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                branch_x - Inches(0.6), branch_y,
                Inches(1.2), Inches(0.8)
            )
            outcome.fill.solid()
            outcome.fill.fore_color.rgb = RGBColor(106, 153, 78)
            outcome.line.width = Pt(2)
            
            tf = outcome.text_frame
            p = tf.paragraphs[0]
            p.text = option
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER

    def generate_lean_canvas(self, prs: Presentation, slide, data: Dict[str, List[str]]) -> None:
        """Generate Lean Canvas (9 blocks)."""
        print("DiagramService: Generating Lean Canvas...")
        
        # Define 9 blocks with positions (different from BMC)
        blocks = [
            # Left column
            ("Problem", Inches(0.5), Inches(1.5), Inches(2.2), Inches(2.2), RGBColor(199, 62, 29)),
            ("Solution", Inches(0.5), Inches(3.8), Inches(2.2), Inches(2.2), RGBColor(106, 153, 78)),
            # Second column
            ("Key Metrics", Inches(2.8), Inches(1.5), Inches(2.2), Inches(2.2), RGBColor(241, 143, 1)),
            ("Unfair Advantage", Inches(2.8), Inches(3.8), Inches(2.2), Inches(2.2), RGBColor(162, 59, 114)),
            # Center column
            ("Unique Value Proposition", Inches(5.1), Inches(1.5), Inches(2.2), Inches(4.5), RGBColor(46, 134, 171)),
            # Fourth column
            ("Channels", Inches(7.4), Inches(1.5), Inches(2.2), Inches(2.2), RGBColor(106, 153, 78)),
            ("Customer Segments", Inches(7.4), Inches(3.8), Inches(2.2), Inches(2.2), RGBColor(241, 143, 1)),
            # Bottom row
            ("Cost Structure", Inches(0.5), Inches(6.1), Inches(4.6), Inches(1.4), RGBColor(150, 150, 150)),
            ("Revenue Streams", Inches(5.1), Inches(6.1), Inches(4.6), Inches(1.4), RGBColor(106, 153, 78))
        ]
        
        for block_name, left, top, width, height, color in blocks:
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left, top, width, height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = color
            shape.line.color.rgb = RGBColor(255, 255, 255)
            shape.line.width = Pt(2)
            
            tf = shape.text_frame
            tf.margin_top = Inches(0.05)
            tf.margin_left = Inches(0.1)
            tf.word_wrap = True
            
            p = tf.paragraphs[0]
            p.text = block_name
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            
            # Add items
            for item in data.get(block_name, [])[:2]:
                p = tf.add_paragraph()
                p.text = f"• {item}"
                p.font.size = Pt(8)
                p.font.color.rgb = RGBColor(255, 255, 255)

    def generate_nine_box(self, prs: Presentation, slide, data: Dict[str, List[str]]) -> None:
        """Generate 9-Box Grid (Talent Assessment)."""
        print("DiagramService: Generating 9-Box Grid...")
        
        # 3x3 grid
        cell_width = Inches(2.8)
        cell_height = Inches(1.8)
        start_left = Inches(1)
        start_top = Inches(1.5)
        
        # Define 9 boxes with labels and colors
        boxes = [
            # Row 1 (High Potential)
            ("High Potential\nLow Performance", 0, 0, RGBColor(241, 143, 1), "Enigma"),
            ("High Potential\nMedium Performance", 1, 0, RGBColor(106, 153, 78), "High Potential"),
            ("High Potential\nHigh Performance", 2, 0, RGBColor(46, 134, 171), "Stars"),
            # Row 2 (Medium Potential)
            ("Medium Potential\nLow Performance", 0, 1, RGBColor(199, 62, 29), "Inconsistent"),
            ("Medium Potential\nMedium Performance", 1, 1, RGBColor(150, 150, 150), "Solid Performers"),
            ("Medium Potential\nHigh Performance", 2, 1, RGBColor(106, 153, 78), "Core Players"),
            # Row 3 (Low Potential)
            ("Low Potential\nLow Performance", 0, 2, RGBColor(199, 62, 29), "Low Performers"),
            ("Low Potential\nMedium Performance", 1, 2, RGBColor(150, 150, 150), "Effective"),
            ("Low Potential\nHigh Performance", 2, 2, RGBColor(106, 153, 78), "Trusted Professionals")
        ]
        
        for label, col, row, color, category in boxes:
            left = start_left + col * cell_width
            top = start_top + row * cell_height
            
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                left, top, cell_width, cell_height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = color
            shape.line.color.rgb = RGBColor(255, 255, 255)
            shape.line.width = Pt(2)
            
            tf = shape.text_frame
            tf.margin_top = Inches(0.05)
            tf.margin_left = Inches(0.1)
            
            p = tf.paragraphs[0]
            p.text = category
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER
            
            # Add employees in this box
            for employee in data.get(category, [])[:2]:
                p = tf.add_paragraph()
                p.text = f"• {employee}"
                p.font.size = Pt(9)
                p.font.color.rgb = RGBColor(255, 255, 255)
                p.alignment = PP_ALIGN.CENTER
        
        # Add axis labels
        # X-axis (Performance)
        perf_label = slide.shapes.add_textbox(
            start_left, start_top + 3 * cell_height + Inches(0.1),
            3 * cell_width, Inches(0.4)
        )
        tf = perf_label.text_frame
        p = tf.paragraphs[0]
        p.text = "Performance →"
        p.font.size = Pt(12)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        
        # Y-axis (Potential)
        pot_label = slide.shapes.add_textbox(
            start_left - Inches(0.8), start_top,
            Inches(0.6), 3 * cell_height
        )
        tf = pot_label.text_frame
        tf.vertical_anchor = 1
        p = tf.paragraphs[0]
        p.text = "Potential ↑"
        p.font.size = Pt(12)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

    def generate_fte_model(self, prs: Presentation, slide, data: Dict[str, Dict]) -> None:
        """Generate FTE (Full-Time Equivalent) Model."""
        print("DiagramService: Generating FTE Model...")
        
        # Bar chart showing Current vs Required FTE
        departments = list(data.keys())[:5]
        bar_width = Inches(1.2)
        bar_spacing = Inches(0.3)
        start_left = Inches(1.5)
        start_top = Inches(5)
        max_height = Inches(3)
        
        # Find max FTE for scaling
        max_fte = max([max(dept.get("current", 0), dept.get("required", 0)) 
                       for dept in data.values()] + [1])
        
        for i, dept_name in enumerate(departments):
            dept_data = data[dept_name]
            current = dept_data.get("current", 0)
            required = dept_data.get("required", 0)
            gap = required - current
            
            x = start_left + i * (2 * bar_width + bar_spacing)
            
            # Current FTE bar
            current_height = max_height * (current / max_fte)
            current_bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                x, start_top - current_height,
                bar_width, current_height
            )
            current_bar.fill.solid()
            current_bar.fill.fore_color.rgb = RGBColor(46, 134, 171)
            current_bar.line.width = Pt(1)
            
            # Current label
            current_text = slide.shapes.add_textbox(
                x, start_top - current_height - Inches(0.3),
                bar_width, Inches(0.25)
            )
            tf = current_text.text_frame
            p = tf.paragraphs[0]
            p.text = f"{current:.1f}"
            p.font.size = Pt(10)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
            
            # Required FTE bar
            required_height = max_height * (required / max_fte)
            required_bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                x + bar_width + Inches(0.1), start_top - required_height,
                bar_width, required_height
            )
            required_bar.fill.solid()
            # Color based on gap
            if gap > 0:
                required_bar.fill.fore_color.rgb = RGBColor(199, 62, 29)  # Red - shortage
            elif gap < 0:
                required_bar.fill.fore_color.rgb = RGBColor(150, 150, 150)  # Grey - excess
            else:
                required_bar.fill.fore_color.rgb = RGBColor(106, 153, 78)  # Green - optimal
            required_bar.line.width = Pt(1)
            
            # Required label
            required_text = slide.shapes.add_textbox(
                x + bar_width + Inches(0.1), start_top - required_height - Inches(0.3),
                bar_width, Inches(0.25)
            )
            tf = required_text.text_frame
            p = tf.paragraphs[0]
            p.text = f"{required:.1f}"
            p.font.size = Pt(10)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
            
            # Department label
            dept_label = slide.shapes.add_textbox(
                x, start_top + Inches(0.1),
                2 * bar_width + Inches(0.1), Inches(0.4)
            )
            tf = dept_label.text_frame
            p = tf.paragraphs[0]
            p.text = dept_name
            p.font.size = Pt(10)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
        
        # Legend
        legend_y = Inches(1)
        legend_items = [
            ("Current FTE", RGBColor(46, 134, 171)),
            ("Required FTE", RGBColor(199, 62, 29))
        ]
        for i, (label, color) in enumerate(legend_items):
            legend_box = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(7) + i * Inches(1.5), legend_y,
                Inches(0.3), Inches(0.3)
            )
            legend_box.fill.solid()
            legend_box.fill.fore_color.rgb = color
            legend_box.line.width = Pt(1)
            
            legend_text = slide.shapes.add_textbox(
                Inches(7.4) + i * Inches(1.5), legend_y,
                Inches(1), Inches(0.3)
            )
            tf = legend_text.text_frame
            p = tf.paragraphs[0]
            p.text = label
            p.font.size = Pt(9)

    def generate_ahp(self, prs: Presentation, slide, data: Dict) -> None:
        """Generate AHP (Analytic Hierarchy Process) Hierarchy."""
        print("DiagramService: Generating AHP...")
        
        # 3-level hierarchy
        goal = data.get("goal", "Decision Goal")
        criteria = data.get("criteria", {})
        alternatives = data.get("alternatives", {})
        
        # Level 1: Goal (top)
        goal_x = Inches(5)
        goal_y = Inches(1.5)
        goal_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            goal_x - Inches(1.5), goal_y,
            Inches(3), Inches(0.8)
        )
        goal_box.fill.solid()
        goal_box.fill.fore_color.rgb = RGBColor(46, 134, 171)
        goal_box.line.width = Pt(3)
        
        tf = goal_box.text_frame
        p = tf.paragraphs[0]
        p.text = goal
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        # Level 2: Criteria (middle)
        criteria_y = Inches(3.5)
        num_criteria = min(len(criteria), 4)
        criteria_spacing = Inches(8) / max(num_criteria, 1)
        
        for i, (criterion, weight) in enumerate(list(criteria.items())[:4]):
            x = Inches(1) + i * criteria_spacing
            
            # Size based on weight
            box_width = Inches(1.2 + weight * 0.8)
            box_height = Inches(0.6)
            
            crit_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x, criteria_y,
                box_width, box_height
            )
            crit_box.fill.solid()
            crit_box.fill.fore_color.rgb = RGBColor(241, 143, 1)
            crit_box.line.width = Pt(2)
            
            tf = crit_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"{criterion}\n{weight*100:.0f}%"
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER
            
            # Connector from goal
            conn = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                goal_x, goal_y + Inches(0.8),
                x + box_width/2, criteria_y
            )
            conn.line.width = Pt(1.5)
        
        # Level 3: Alternatives (bottom)
        alt_y = Inches(5.5)
        num_alt = min(len(alternatives), 3)
        alt_spacing = Inches(8) / max(num_alt, 1)
        
        for i, (alternative, score) in enumerate(list(alternatives.items())[:3]):
            x = Inches(1.5) + i * alt_spacing
            
            # Size based on score
            box_width = Inches(1.5 + score * 0.8)
            box_height = Inches(0.6)
            
            alt_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x, alt_y,
                box_width, box_height
            )
            alt_box.fill.solid()
            alt_box.fill.fore_color.rgb = RGBColor(106, 153, 78)
            alt_box.line.width = Pt(2)
            
            tf = alt_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"{alternative}\n{score*100:.0f}%"
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER

    def generate_ocai(self, prs: Presentation, slide, data: Dict[str, float]) -> None:
        """Generate OCAI (Organizational Culture Assessment) - 4 quadrants."""
        print("DiagramService: Generating OCAI...")
        
        # 4 culture types
        cultures = ["Clan", "Adhocracy", "Market", "Hierarchy"]
        colors = [
            RGBColor(106, 153, 78),   # Clan - Green
            RGBColor(241, 143, 1),    # Adhocracy - Orange
            RGBColor(199, 62, 29),    # Market - Red
            RGBColor(46, 134, 171)    # Hierarchy - Blue
        ]
        
        # Center point
        center_x = Inches(5)
        center_y = Inches(4)
        
        # Draw quadrants
        quad_size = Inches(3)
        positions = [
            (center_x - quad_size, center_y - quad_size),  # Top-left (Clan)
            (center_x, center_y - quad_size),              # Top-right (Adhocracy)
            (center_x, center_y),                          # Bottom-right (Market)
            (center_x - quad_size, center_y)               # Bottom-left (Hierarchy)
        ]
        
        for i, (culture, color) in enumerate(zip(cultures, colors)):
            left, top = positions[i]
            score = data.get(culture, 0.25)  # Default 25%
            
            quad = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                left, top, quad_size, quad_size
            )
            quad.fill.solid()
            quad.fill.fore_color.rgb = color
            quad.fill.transparency = 0.3
            quad.line.color.rgb = RGBColor(255, 255, 255)
            quad.line.width = Pt(2)
            
            # Label
            label_box = slide.shapes.add_textbox(
                left + Inches(0.2), top + Inches(0.2),
                Inches(1.5), Inches(0.8)
            )
            tf = label_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"{culture}\n{score*100:.0f}%"
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = RGBColor(50, 50, 50)
        
        
        # Center circle
        center_circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            center_x - Inches(0.5), center_y - Inches(0.5),
            Inches(1), Inches(1)
        )
        center_circle.fill.solid()
        center_circle.fill.fore_color.rgb = RGBColor(255, 255, 255)
        center_circle.line.width = Pt(3)
        
        tf = center_circle.text_frame
        p = tf.paragraphs[0]
        p.text = "OCAI"
        p.font.size = Pt(12)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

    def generate_futures_wheel(self, prs: Presentation, slide, data: Dict[str, List[str]]) -> None:
        """Generate Futures Wheel (Central event + consequences)."""
        print("DiagramService: Generating Futures Wheel...")
        
        # Central Event
        center_x = Inches(5)
        center_y = Inches(3.75)
        
        center_circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            center_x - Inches(0.8), center_y - Inches(0.8),
            Inches(1.6), Inches(1.6)
        )
        center_circle.fill.solid()
        center_circle.fill.fore_color.rgb = RGBColor(46, 134, 171)
        center_circle.line.width = Pt(2)
        
        tf = center_circle.text_frame
        p = tf.paragraphs[0]
        p.text = data.get("event", "Central Event")
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        # 1st Order Consequences (surrounding)
        import math
        consequences = data.get("consequences", [])[:6]
        radius = Inches(2.5)
        
        for i, text in enumerate(consequences):
            angle = (i * 360 / len(consequences)) * math.pi / 180
            x = center_x + radius * math.cos(angle)
            y = center_y + radius * math.sin(angle)
            
            # Connector line
            line = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                center_x, center_y,
                x, y
            )
            line.line.width = Pt(1.5)
            line.line.color.rgb = RGBColor(100, 100, 100)
            
            # Consequence Bubble
            bubble = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                x - Inches(0.6), y - Inches(0.6),
                Inches(1.2), Inches(1.2)
            )
            bubble.fill.solid()
            bubble.fill.fore_color.rgb = RGBColor(241, 143, 1)
            bubble.line.width = Pt(1)
            
            tf = bubble.text_frame
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER

    def generate_kano_model(self, prs: Presentation, slide, data: Dict[str, List[str]]) -> None:
        """Generate Kano Model (Satisfaction vs Implementation)."""
        print("DiagramService: Generating Kano Model...")
        
        # Axes
        start_x = Inches(1.5)
        end_x = Inches(8.5)
        start_y = Inches(5.5)
        end_y = Inches(1.5)
        
        # X-axis
        x_axis = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            start_x, start_y, end_x, start_y
        )
        x_axis.line.width = Pt(2)
        x_axis.line.color.rgb = RGBColor(50, 50, 50)
        
        # Y-axis
        y_axis = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            start_x, start_y, start_x, end_y
        )
        y_axis.line.width = Pt(2)
        y_axis.line.color.rgb = RGBColor(50, 50, 50)
        
        # Labels
        x_label = slide.shapes.add_textbox(end_x, start_y, Inches(1.5), Inches(0.5))
        x_label.text_frame.text = "Functionality"
        
        y_label = slide.shapes.add_textbox(start_x, end_y - Inches(0.5), Inches(1.5), Inches(0.5))
        y_label.text_frame.text = "Satisfaction"
        
        # Curves (Simplified as curved lines using shapes or just text labels for regions)
        # Since drawing bezier curves is complex in pptx without custom geometry, 
        # we will use curved arrows or approximations.
        
        # Basic Needs (Curve down-right)
        basic_curve = slide.shapes.add_shape(
            MSO_SHAPE.ARC,
            start_x, start_y - Inches(1.5), Inches(3), Inches(1.5)
        )
        basic_curve.line.color.rgb = RGBColor(199, 62, 29) # Red
        basic_curve.line.width = Pt(3)
        basic_curve.rotation = 180
        
        slide.shapes.add_textbox(start_x + Inches(2), start_y - Inches(0.5), Inches(1), Inches(0.5)).text_frame.text = "Basic"
        
        # Performance (Straight diagonal)
        perf_line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            start_x, start_y,
            start_x + Inches(4), start_y - Inches(2.5)
        )
        perf_line.line.color.rgb = RGBColor(46, 134, 171) # Blue
        perf_line.line.width = Pt(3)
        
        slide.shapes.add_textbox(start_x + Inches(2.5), start_y - Inches(1.8), Inches(1.5), Inches(0.5)).text_frame.text = "Performance"
        
        # Delighters (Curve up-left)
        delight_curve = slide.shapes.add_shape(
            MSO_SHAPE.ARC,
            start_x, start_y - Inches(4), Inches(3), Inches(1.5)
        )
        delight_curve.line.color.rgb = RGBColor(106, 153, 78) # Green
        delight_curve.line.width = Pt(3)
        delight_curve.rotation = 270 # Adjust rotation to look like top curve
        
        # Better approximation for Delighters: Arc starting from Y-axis going up and right
        # Actually, let's just place text labels for features in approximate positions
        
        # Plot Features
        for feature in data.get("Basic", []):
            # Bottom right area
            box = slide.shapes.add_textbox(start_x + Inches(5), start_y - Inches(0.5), Inches(1.5), Inches(0.3))
            box.text_frame.text = f"• {feature}"
            box.text_frame.paragraphs[0].font.size = Pt(9)
            
        for feature in data.get("Performance", []):
            # Middle diagonal area
            box = slide.shapes.add_textbox(start_x + Inches(3.5), start_y - Inches(2.5), Inches(1.5), Inches(0.3))
            box.text_frame.text = f"• {feature}"
            box.text_frame.paragraphs[0].font.size = Pt(9)
            
        for feature in data.get("Delighters", []):
            # Top left area
            box = slide.shapes.add_textbox(start_x + Inches(1), start_y - Inches(3.5), Inches(1.5), Inches(0.3))
            box.text_frame.text = f"• {feature}"
            box.text_frame.paragraphs[0].font.size = Pt(9)

