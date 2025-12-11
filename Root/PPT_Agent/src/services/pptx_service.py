from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from src.models.schema import PresentationSchema
from src.services.diagram_service import DiagramService
from src.services.data_ingestion_service import DataIngestionService
from src.services.job_data_extractor import JobDataExtractor
import os

class PPTXService:
    def __init__(self, output_dir: str = "output"):
        self.output_dir = output_dir
        os.makedirs(self.output_dir, exist_ok=True)
        self.diagram_service = DiagramService()
        self.data_ingestion = DataIngestionService()
        self.extractor = JobDataExtractor()

    def create_presentation(self, schema: PresentationSchema, filename: str = "presentation.pptx", design_schema: list = None) -> str:
        prs = Presentation()

        # Content Slides
        bullet_slide_layout = prs.slide_layouts[1]

        for i, slide_content in enumerate(schema.slides):
            # Determine layout (default to bullet if not specified)
            layout_idx = 1
            style = {}
            if design_schema and i < len(design_schema):
                layout_idx = design_schema[i].get("layout_index", 1)
                style = design_schema[i]

            slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
            shapes = slide.shapes

            # Apply Background Color
            if "background_color" in style:
                bg = slide.background
                fill = bg.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(*style["background_color"])

            # Title
            title_shape = shapes.title
            title_shape.text = slide_content.title
            if "title_color" in style:
                title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(*style["title_color"])
                title_shape.text_frame.paragraphs[0].font.name = style.get("font_family", "Calibri")

            # Body / Bullet Points
            if slide_content.bullet_points:
                body_shape = shapes.placeholders[1]
                tf = body_shape.text_frame
                tf.text = slide_content.bullet_points[0]
                
                # Apply style to first point
                if "body_color" in style:
                    tf.paragraphs[0].font.color.rgb = RGBColor(*style["body_color"])
                    tf.paragraphs[0].font.name = style.get("font_family", "Calibri")

                for point in slide_content.bullet_points[1:]:
                    p = tf.add_paragraph()
                    p.text = point
                    p.level = 1
                    if "body_color" in style:
                        p.font.color.rgb = RGBColor(*style["body_color"])
                        p.font.name = style.get("font_family", "Calibri")

            # Diagram (takes highest precedence)
            if "diagram_type" in style:
                diagram_type = style["diagram_type"]
                
                # Resolve data
                diagram_data = style.get("diagram_data", {})
                
                # Hybrid Data Strategy: Check for data_source if diagram_data is empty
                if not diagram_data and "data_source" in style:
                    source = style["data_source"]
                    source_type = source.get("type")
                    
                    if source_type == "db":
                        # Fetch from DB
                        inst_id = source.get("institution_id")
                        if inst_id:
                            diagram_data = self.extractor.get_strategic_analysis(inst_id, diagram_type.upper())
                    elif source_type == "direct":
                        # Ingest direct content
                        content = source.get("content")
                        diagram_data = self.data_ingestion.ingest_data(content, "dict")
                
                if diagram_type == "swot":
                    self.diagram_service.generate_swot(prs, slide, diagram_data)
                elif diagram_type == "workflow":
                    self.diagram_service.generate_workflow(prs, slide, diagram_data.get("steps", []))
                elif diagram_type == "4p":
                    self.diagram_service.generate_four_p(prs, slide, diagram_data)
                elif diagram_type == "bcg":
                    self.diagram_service.generate_bcg_matrix(prs, slide, diagram_data)
                elif diagram_type == "eisenhower":
                    self.diagram_service.generate_eisenhower_matrix(prs, slide, diagram_data)
                elif diagram_type == "porters":
                    self.diagram_service.generate_porters_five_forces(prs, slide, diagram_data)
                elif diagram_type == "ge_matrix":
                    self.diagram_service.generate_ge_matrix(prs, slide, diagram_data)
                elif diagram_type == "scenario":
                    self.diagram_service.generate_scenario_planning(prs, slide, diagram_data)
                elif diagram_type == "value_chain":
                    self.diagram_service.generate_value_chain(prs, slide, diagram_data)
                elif diagram_type == "bmc":
                    self.diagram_service.generate_business_model_canvas(prs, slide, diagram_data)
                elif diagram_type == "mckinsey_7s":
                    self.diagram_service.generate_mckinsey_7s(prs, slide, diagram_data)
                elif diagram_type == "ansoff":
                    self.diagram_service.generate_ansoff_matrix(prs, slide, diagram_data)
                elif diagram_type == "stakeholder":
                    self.diagram_service.generate_stakeholder_mapping(prs, slide, diagram_data)
                elif diagram_type == "3c":
                    self.diagram_service.generate_3c_analysis(prs, slide, diagram_data)
                elif diagram_type == "impact_effort":
                    self.diagram_service.generate_impact_effort_matrix(prs, slide, diagram_data)
                elif diagram_type == "gantt":
                    self.diagram_service.generate_gantt_chart(prs, slide, diagram_data)
                elif diagram_type == "fishbone":
                    self.diagram_service.generate_fishbone(prs, slide, diagram_data)
                elif diagram_type == "org_chart":
                    self.diagram_service.generate_org_chart(prs, slide, diagram_data)
                elif diagram_type == "raci":
                    self.diagram_service.generate_raci_matrix(prs, slide, diagram_data)
                elif diagram_type == "decision_tree":
                    self.diagram_service.generate_decision_tree(prs, slide, diagram_data)
                elif diagram_type == "lean_canvas":
                    self.diagram_service.generate_lean_canvas(prs, slide, diagram_data)
                elif diagram_type == "nine_box":
                    self.diagram_service.generate_nine_box(prs, slide, diagram_data)
                elif diagram_type == "fte_model":
                    self.diagram_service.generate_fte_model(prs, slide, diagram_data)
                elif diagram_type == "ahp":
                    self.diagram_service.generate_ahp(prs, slide, diagram_data)
                elif diagram_type == "ocai":
                    self.diagram_service.generate_ocai(prs, slide, diagram_data)
                elif diagram_type == "futures_wheel":
                    self.diagram_service.generate_futures_wheel(prs, slide, diagram_data)
                elif diagram_type == "kano_model":
                    self.diagram_service.generate_kano_model(prs, slide, diagram_data)
            
            # Chart (if no diagram)
            elif "chart_path" in style and style["chart_path"]:
                left = Inches(5.5)
                top = Inches(1.5)
                height = Inches(4)
                slide.shapes.add_picture(style["chart_path"], left, top, height=height)
            
            # Image (if no diagram or chart)
            elif slide_content.image_description and "image_path" in style and style["image_path"]:
                # Add image to the right side (simple placement)
                left = Inches(5.5)
                top = Inches(2)
                height = Inches(3.5)
                slide.shapes.add_picture(style["image_path"], left, top, height=height)

            if slide_content.speaker_notes:
                slide.notes_slide.notes_text_frame.text = slide_content.speaker_notes

        output_path = os.path.join(self.output_dir, filename)
        prs.save(output_path)
        return output_path
