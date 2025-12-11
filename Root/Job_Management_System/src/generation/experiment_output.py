import os
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pathlib import Path

OUTPUT_DIR = Path("data/output")
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
# We now default to the complex strategy JSON
INPUT_JSON = OUTPUT_DIR / "complex_strategy.json"

def load_data():
    if not INPUT_JSON.exists():
        # Fallback to original if complex doesn't exist
        fallback = OUTPUT_DIR / "slide_analysis.json"
        if fallback.exists():
            return json.load(open(fallback, "r", encoding="utf-8"))
        return None
    with open(INPUT_JSON, "r", encoding="utf-8") as f:
        return json.load(f)

def generate_html(data):
    if not data: return

    # --- CSS Styles ---
    css = """
    @import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;600;700&display=swap');

    body {
        font-family: 'Inter', sans-serif;
        display: flex;
        justify-content: center;
        align-items: center;
        min-height: 100vh;
        background: #e0e5ec;
        margin: 0;
    }

    .slide-container {
        width: 960px;
        height: 540px;
        background: white;
        box-shadow: 0 20px 50px rgba(0,0,0,0.1);
        position: relative;
        overflow: hidden;
        display: flex;
        flex-direction: column;
        padding: 40px 50px;
        box-sizing: border-box;
    }

    /* Header */
    .header {
        border-bottom: 2px solid #003366;
        padding-bottom: 10px;
        margin-bottom: 20px;
        display: flex;
        justify-content: space-between;
        align-items: flex-end;
    }
    .header h1 {
        font-size: 24px;
        color: #003366;
        margin: 0;
        font-weight: 700;
        text-transform: uppercase;
        letter-spacing: 0.5px;
    }
    .governing-message {
        font-size: 14px;
        color: #444;
        margin-bottom: 20px;
        padding: 10px 15px;
        background: #f1f3f5;
        border-left: 4px solid #003366;
        font-weight: 500;
    }

    /* Layout Grid */
    .content-area {
        display: flex;
        flex-direction: column;
        gap: 15px;
        flex: 1;
    }

    /* Recursive Block Rendering */
    .block-container {
        padding: 10px;
        border-radius: 4px;
        display: flex;
        flex-direction: column;
        gap: 8px;
    }
    .block-container h4 {
        margin: 0 0 8px 0;
        font-size: 13px;
        color: #003366;
        border-bottom: 1px solid #ccc;
        padding-bottom: 4px;
    }

    .block-text { font-size: 11px; line-height: 1.4; color: #333; }
    .block-kpi { 
        background: white; 
        padding: 8px; 
        border-radius: 4px; 
        text-align: center; 
        box-shadow: 0 2px 4px rgba(0,0,0,0.05);
        border: 1px solid #eee;
    }

    /* Strategy House Layouts */
    .layout-roof {
        width: 100%;
        margin-bottom: 10px;
    }
    .layout-columns-3 {
        display: grid;
        grid-template-columns: 1fr 1fr 1fr;
        gap: 15px;
    }
    .layout-full-width {
        width: 100%;
    }
    
    /* Process Flow */
    .layout-process-flow {
        display: flex;
        gap: 10px;
        margin-top: auto; /* Push to bottom of content area */
    }
    .chevron {
        flex: 1;
        position: relative;
        padding: 10px 20px 10px 30px;
        font-size: 11px;
        font-weight: 600;
        display: flex;
        align-items: center;
        justify-content: center;
        clip-path: polygon(0% 0%, 90% 0%, 100% 50%, 90% 100%, 0% 100%, 10% 50%);
        margin-right: -15px; /* Overlap */
        z-index: 1;
    }
    .chevron:first-child {
        clip-path: polygon(0% 0%, 90% 0%, 100% 50%, 90% 100%, 0% 100%);
        padding-left: 20px;
    }

    /* Footer */
    .footer {
        position: absolute;
        bottom: 10px;
        left: 50px;
        right: 50px;
        border-top: 1px solid #eee;
        padding-top: 5px;
        font-size: 9px;
        color: #999;
        display: flex;
        justify-content: space-between;
    }
    """

    # --- Recursive Renderer ---
    def render_block(block):
        b_type = block.get('block_type', 'text')
        content = block.get('content', '')
        style = block.get('style', {})
        
        # Style String Construction
        s_str = ""
        if style.get('background_hex'): s_str += f"background-color: {style['background_hex']}; "
        if style.get('color_hex'): s_str += f"color: {style['color_hex']}; "
        if style.get('font_size'): s_str += f"font-size: {style['font_size']}px; "
        if style.get('is_bold'): s_str += "font-weight: bold; "
        if style.get('border'): s_str += f"border: {style['border']}; "
        if style.get('padding'): s_str += f"padding: {style['padding']}; "
        if style.get('text_align'): s_str += f"text-align: {style['text_align']}; "

        if b_type == 'container':
            children_html = "".join([render_block(b) for b in block.get('blocks', [])])
            title_html = f"<h4>{block.get('title')}</h4>" if block.get('title') else ""
            return f'<div class="block-container" style="{s_str}">{title_html}{children_html}</div>'
        
        elif b_type == 'chevron':
            return f'<div class="chevron" style="{s_str}">{content}</div>'
        
        elif b_type == 'kpi_card':
            return f'<div class="block-kpi" style="{s_str}">{content}</div>'
            
        elif b_type == 'image':
             return f'<div style="background:#eee; height:50px; display:flex; align-items:center; justify-content:center; font-size:10px; color:#888; border:1px dashed #ccc;">IMG: {content}</div>'

        else: # text
            return f'<div class="block-text" style="{s_str}">{content}</div>'

    # --- Section Renderer ---
    sections_html = ""
    for section in data.get('sections', []):
        l_type = section.get('layout_type', 'full_width')
        
        # Map layout type to CSS class
        css_class = f"layout-{l_type.replace('_', '-')}"
        
        # Render children
        blocks_html = "".join([render_block(b) for b in section.get('blocks', [])])
        
        # Section Title (optional, usually handled by blocks or implicit)
        # For Strategy House, we might want a section wrapper
        sections_html += f"""
        <div class="{css_class}">
            {blocks_html}
        </div>
        """

    html_body = f"""
    <div class="slide-container">
        <div class="header">
            <h1>{data.get('title', 'Untitled')}</h1>
            <div style="font-size: 10px; color: #666; align-self: center;">STRATEGY 2030</div>
        </div>
        <div class="governing-message">{data.get('governing_message', '')}</div>
        <div class="content-area">
            {sections_html}
        </div>
        <div class="footer">
            <div>{data.get('footer', '')}</div>
            <div>Page {data.get('slide_number', 1)}</div>
        </div>
    </div>
    """

    full_html = f"""
    <!DOCTYPE html>
    <html>
    <head>
        <style>{css}</style>
    </head>
    <body>
        {html_body}
    </body>
    </html>
    """

    output_path = OUTPUT_DIR / "test_slide.html"
    with open(output_path, "w", encoding="utf-8") as f:
        f.write(full_html)
    print(f"Generated HTML: {output_path}")

def render_slide_to_pptx(prs, data):
    """
    Renders a single slide into the provided Presentation object.
    Handles both 'Strategy House' custom layouts and generic Vision Agent outputs.
    Includes heuristics to infer shape-based layouts from content patterns.
    """
    if not data: return

    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    import re

    # Use a blank layout (6)
    slide_layout = prs.slide_layouts[6] 
    slide = prs.slides.add_slide(slide_layout)
    
    # Dimensions (16:9 default is 10 x 5.625 inches)
    SLIDE_W = prs.slide_width
    SLIDE_H = prs.slide_height
    MARGIN = Inches(0.5)
    CONTENT_W = SLIDE_W - (MARGIN * 2)
    CONTENT_H = SLIDE_H - (MARGIN * 2)

    # --- Helper Functions ---
    def hex_to_rgb(hex_str):
        if not hex_str or not isinstance(hex_str, str) or not hex_str.startswith('#'): return RGBColor(0, 0, 0)
        hex_str = hex_str.lstrip('#')
        try:
            return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))
        except:
            return RGBColor(0, 0, 0)

    def add_text_box(slide, text, left, top, width, height, font_size=10, is_bold=False, align=PP_ALIGN.LEFT, color=None, bg_color=None):
        if text is None: text = ""
        shape = slide.shapes.add_textbox(left, top, width, height)
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = str(text)
        p.font.size = Pt(font_size)
        p.font.bold = is_bold
        p.font.name = "Malgun Gothic" # Korean Font Support
        p.alignment = align
        if color:
            p.font.color.rgb = hex_to_rgb(color)
        return shape

    def add_shape_with_text(slide, shape_type, text, left, top, width, height, bg_color=None, font_color="#000000", font_size=10, is_bold=False):
        if text is None: text = ""
        shape = slide.shapes.add_shape(shape_type, left, top, width, height)
        
        if bg_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = hex_to_rgb(bg_color)
            shape.line.color.rgb = hex_to_rgb(bg_color) # Match border to fill for clean look
        
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = str(text)
        p.font.size = Pt(font_size)
        p.font.bold = is_bold
        p.font.name = "Malgun Gothic"
        p.font.color.rgb = hex_to_rgb(font_color)
        p.alignment = PP_ALIGN.CENTER
        return shape

    def detect_layout_heuristic(blocks):
        """
        Analyzes blocks to guess the intended layout if not specified.
        Returns: 'process_flow', 'columns_3', 'list', or 'generic'
        """
        if not blocks: return 'generic'
        
        # Check for Process Flow (Step 1, Phase A, etc.)
        process_keywords = [r'^step\s?\d', r'^phase\s?\d', r'^\d+\.?\s*단계', r'^\d+\.?\s*step']
        match_count = 0
        for b in blocks:
            content = str(b.get('content', '')).lower()
            if any(re.match(p, content) for p in process_keywords):
                match_count += 1
        
        if match_count >= 2 and len(blocks) <= 8: # If multiple steps found
            return 'process_flow'

        # Check for 3-Column (if exactly 3 blocks and short titles)
        if len(blocks) == 3:
            return 'columns_3'

        return 'generic'

    # --- Header ---
    title = data.get('title', 'Untitled')
    add_text_box(slide, title, MARGIN, Inches(0.2), CONTENT_W, Inches(0.5), font_size=24, is_bold=True, color="#003366")
    
    gov_msg = data.get('governing_message', '')
    if gov_msg:
        add_text_box(slide, gov_msg, MARGIN, Inches(0.8), CONTENT_W, Inches(0.6), font_size=14, color="#555555")

    # --- Content Layout Engine ---
    current_y = Inches(1.6)
    
    # Normalize Sections: Check top-level OR content_data
    sections = data.get('sections')
    if not sections:
        sections = data.get('content_data', {}).get('sections')
    
    if not sections: sections = []

    # Check Grid Type
    grid_type = data.get('layout', {}).get('grid_type', '1-column')
    
    # Check Framework Type (Global Override)
    framework_type = data.get('content_data', {}).get('framework_type')

    # Column Trackers for Generic Layouts
    col_gap = Inches(0.2)
    col_2_w = (CONTENT_W - col_gap) / 2
    left_col_y = current_y
    right_col_y = current_y

    for section in sections:
        if not section: continue
        
        blocks = section.get('blocks') or []
        
        # 1. Explicit Layout
        l_type = section.get('layout_type')
        
        # 2. Inferred Layout from Framework Type (if not explicit)
        if not l_type and framework_type:
            if framework_type == 'Process': l_type = 'process_flow'
            elif framework_type == 'Pyramid': l_type = 'roof' # Approximation
        
        # 3. Heuristic Layout (if still unknown)
        if not l_type:
            l_type = detect_layout_heuristic(blocks)

        # --- Specific Layouts (Strategy House, etc.) ---
        if l_type == 'roof':
            h = Inches(0.8)
            for block in blocks:
                add_shape_with_text(slide, MSO_SHAPE.ISOSCELES_TRIANGLE, block.get('content'), 
                                    MARGIN, current_y, CONTENT_W, h, 
                                    bg_color=block.get('style', {}).get('background_hex', '#003366'),
                                    font_color=block.get('style', {}).get('color_hex', '#FFFFFF'),
                                    font_size=16, is_bold=True)
            current_y += h + Inches(0.1)

        elif l_type == 'columns_3':
            col_w = (CONTENT_W - Inches(0.2)) / 3
            h = Inches(2.5)
            for i, block in enumerate(blocks):
                x = MARGIN + (col_w + Inches(0.1)) * i
                container = add_shape_with_text(slide, MSO_SHAPE.RECTANGLE, "", x, current_y, col_w, h, bg_color="#f8f9fa")
                container.line.color.rgb = hex_to_rgb("#dee2e6")
                
                title_h = Inches(0.4)
                add_text_box(slide, block.get('title', ''), x + Inches(0.1), current_y + Inches(0.1), col_w - Inches(0.2), title_h, 
                             font_size=11, is_bold=True, align=PP_ALIGN.CENTER, color="#003366")
                
                inner_y = current_y + title_h + Inches(0.1)
                for child in block.get('blocks', []):
                    b_type = child.get('block_type', 'text')
                    content = child.get('content', '')
                    if b_type == 'kpi_card':
                        kpi_h = Inches(0.4)
                        add_shape_with_text(slide, MSO_SHAPE.ROUNDED_RECTANGLE, content, 
                                            x + Inches(0.2), inner_y + Inches(0.1), col_w - Inches(0.4), kpi_h,
                                            bg_color="FFFFFF", font_color=child.get('style', {}).get('color_hex', '#000000'),
                                            font_size=10, is_bold=True)
                        inner_y += kpi_h + Inches(0.1)
                    else:
                        add_text_box(slide, content, x + Inches(0.1), inner_y, col_w - Inches(0.2), Inches(0.25), font_size=9)
                        inner_y += Inches(0.25)
            current_y += h + Inches(0.1)

        elif l_type == 'process_flow':
            count = len(blocks)
            # Filter out empty blocks just in case
            valid_blocks = [b for b in blocks if b.get('content')]
            count = len(valid_blocks)
            if count == 0: continue

            w = (CONTENT_W + Inches(0.2)) / count 
            h = Inches(0.6) # Slightly taller for process
            
            for i, block in enumerate(valid_blocks):
                x = MARGIN + (w - Inches(0.15)) * i 
                shape_type = MSO_SHAPE.CHEVRON
                if i == 0: shape_type = MSO_SHAPE.PENTAGON 
                
                # Alternate colors for process steps if no color specified
                default_bg = '#0056b3' if i % 2 == 0 else '#007bff'
                
                add_shape_with_text(slide, shape_type, block.get('content'), 
                                    x, current_y, w, h, 
                                    bg_color=block.get('style', {}).get('background_hex', default_bg),
                                    font_color="#FFFFFF", font_size=10, is_bold=True)
            current_y += h + Inches(0.2)

        # --- Generic / Fallback Layout Logic ---
        else:
            # Handle generic blocks based on position or grid type
            section_title = section.get('title')
            if section_title:
                # If we are in a 2-column mode, we need to decide where to put the title
                # For now, let's put it full width if it's a section title
                add_text_box(slide, section_title, MARGIN, max(left_col_y, right_col_y), CONTENT_W, Inches(0.3), font_size=12, is_bold=True, color="#003366")
                left_col_y += Inches(0.4)
                right_col_y += Inches(0.4)

            for block in blocks:
                content = block.get('content', '')
                pos = block.get('position') or 'top-left'
                
                # Determine X, Y, W
                x, y, w = MARGIN, left_col_y, CONTENT_W
                
                if grid_type == '2-column':
                    if 'right' in pos:
                        x = MARGIN + col_2_w + col_gap
                        y = right_col_y
                        w = col_2_w
                    else: # left
                        x = MARGIN
                        y = left_col_y
                        w = col_2_w
                else:
                    y = left_col_y # Default to single column stack

                # Render Block
                # Estimate height based on text length (very rough)
                lines = len(str(content)) // 50 + 1
                h = Inches(0.25 * lines)
                
                if block.get('block_type') == 'table':
                    # Render table as text box for now (simplified)
                    add_text_box(slide, f"[TABLE]\n{content}", x, y, w, h, font_size=9, bg_color=None)
                else:
                    add_text_box(slide, content, x, y, w, h, font_size=10)

                # Update Y trackers
                if grid_type == '2-column':
                    if 'right' in pos: right_col_y += h + Inches(0.1)
                    else: left_col_y += h + Inches(0.1)
                else:
                    left_col_y += h + Inches(0.1)
                    right_col_y = left_col_y

    # --- Footer ---
    footer_text = data.get('footer', '')
    add_text_box(slide, f"{footer_text} | Page {data.get('slide_number', 1)}", MARGIN, SLIDE_H - Inches(0.4), CONTENT_W, Inches(0.3), font_size=9, color="#999999", align=PP_ALIGN.RIGHT)


def generate_pptx(data):
    if not data: return
    from pptx import Presentation
    
    prs = Presentation()
    render_slide_to_pptx(prs, data)
    
    output_path = OUTPUT_DIR / "korean_strategy.pptx"
    prs.save(output_path)
    print(f"Generated PPTX: {output_path}")

if __name__ == "__main__":
    print("Generating outputs...")
    data = load_data()
    if data:
        generate_html(data)
        generate_pptx(data)
    print("Done.")
